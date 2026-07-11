import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    BG_COLOR = RGBColor(7, 10, 19)         # Deep Space / Navy (Dark Background)
    CARD_BG = RGBColor(18, 25, 41)          # Sleek card background
    TEXT_MAIN = RGBColor(255, 255, 255)     # Crisp white
    TEXT_MUTED = RGBColor(160, 174, 192)    # Muted grey-blue
    ACCENT_GOLD = RGBColor(245, 158, 11)    # Lumens / Gold accent
    ACCENT_CYAN = RGBColor(6, 182, 212)     # Trust / Stellar Cyan accent
    CARD_BORDER = RGBColor(30, 41, 59)      # Slate border

    blank_layout = prs.slide_layouts[6]

    # Helper function to add a consistent background
    def add_background(slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background() # No border
        return bg

    # Helper function to add slide title
    def add_slide_header(slide, title_text, category_text=None):
        if category_text:
            cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
            cat_tf = cat_box.text_frame
            cat_tf.word_wrap = True
            p0 = cat_tf.paragraphs[0]
            p0.text = category_text.upper()
            p0.font.name = 'Segoe UI'
            p0.font.size = Pt(10)
            p0.font.bold = True
            p0.font.color.rgb = ACCENT_CYAN

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p.alignment = PP_ALIGN.LEFT

    # Helper function to add a card shape
    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # Helper to add text to a card
    def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_MAIN, align=PP_ALIGN.LEFT, font_name='Segoe UI'):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return box

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark Space Premium)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_background(slide1)

    # Decorative large abstract circle representing a "Lumen" or star
    circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(1.5), Inches(4.5), Inches(4.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(12, 20, 39)
    circle.line.color.rgb = ACCENT_CYAN
    circle.line.width = Pt(1)

    # Sub-glow circle
    circle_inner = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(2.5), Inches(2.5), Inches(2.5))
    circle_inner.fill.solid()
    circle_inner.fill.fore_color.rgb = RGBColor(16, 32, 60)
    circle_inner.line.color.rgb = ACCENT_GOLD
    circle_inner.line.width = Pt(1.5)

    # Title & Subtitle block
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(7.5), Inches(4.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = "LUMENOVA"
    p_title.font.name = 'Trebuchet MS'
    p_title.font.size = Pt(64)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_MAIN
    
    p_tag = tf.add_paragraph()
    p_tag.text = "Trustless Crowdfunding & Milestone-Based Escrow"
    p_tag.font.name = 'Segoe UI'
    p_tag.font.size = Pt(22)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_GOLD
    p_tag.space_before = Pt(15)

    p_desc = tf.add_paragraph()
    p_desc.text = "Empowering transparent project funding on Stellar Soroban through community governance, weighted on-chain voting, and soulbound badges."
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.space_before = Pt(10)

    # Stellar badge
    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.5), Inches(3.2), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(15, 23, 42)
    badge.line.color.rgb = ACCENT_CYAN
    badge_tf = badge.text_frame
    badge_tf.word_wrap = True
    badge_p = badge_tf.paragraphs[0]
    badge_p.text = "✦ Built for Stellar Builder Challenge L5"
    badge_p.font.name = 'Segoe UI'
    badge_p.font.size = Pt(11)
    badge_p.font.bold = True
    badge_p.font.color.rgb = ACCENT_CYAN
    badge_p.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 2: Problem Statement
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_background(slide2)
    add_slide_header(slide2, "The Critical Trust Gap in Crowdfunding", "Problem Statement")

    # Column 1: Traditional Crowdfunding
    add_card(slide2, Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.5))
    add_text_box(slide2, Inches(1.3), Inches(2.1), Inches(4.6), Inches(0.5), "Traditional Escrow Model", 18, True, ACCENT_GOLD)
    
    box_left = slide2.shapes.add_textbox(Inches(1.3), Inches(2.7), Inches(4.6), Inches(3.2))
    tf_left = box_left.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "• 100% Upfront Disbursement:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_left.add_paragraph()
    p.text = "Creators receive all capital upon campaign success, leaving backers with zero leverage."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)
    
    p = tf_left.add_paragraph()
    p.text = "• No Ongoing Verification:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_left.add_paragraph()
    p.text = "No mechanism exists to verify if project milestones are actually reached before funds are spent."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = tf_left.add_paragraph()
    p.text = "• High Risk of Capital Loss:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_left.add_paragraph()
    p.text = "If a campaign is abandoned or delayed, donors have no recourse and cannot retrieve their funds."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED

    # Middle Separation Visual (Gap indicator)
    middle_arrow = slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(3.5), Inches(0.4), Inches(0.8))
    middle_arrow.fill.solid()
    middle_arrow.fill.fore_color.rgb = RGBColor(220, 38, 38) # Red alert
    middle_arrow.line.fill.background()

    # Column 2: Web3 / Standard Smart Escrows
    add_card(slide2, Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.5))
    add_text_box(slide2, Inches(7.4), Inches(2.1), Inches(4.6), Inches(0.5), "The Backer's Risk Profile", 18, True, RGBColor(239, 68, 68))
    
    box_right = slide2.shapes.add_textbox(Inches(7.4), Inches(2.7), Inches(4.6), Inches(3.2))
    tf_right = box_right.text_frame
    tf_right.word_wrap = True

    p = tf_right.paragraphs[0]
    p.text = "• Blind Trust Dependency:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_right.add_paragraph()
    p.text = "Backers are forced to trust social media updates or written reports, which can be easily fabricated."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)
    
    p = tf_right.add_paragraph()
    p.text = "• Platform Centralization:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_right.add_paragraph()
    p.text = "Web2 platforms take massive cuts (5-10%) and can freeze accounts arbitrarily without community consensus."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = tf_right.add_paragraph()
    p.text = "• Disincentivized Backers:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_right.add_paragraph()
    p.text = "Lack of alignment of incentives leads to low conversion rates, trust fatigue, and fewer repeating donors."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 3: Solution (The Core Framework)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_background(slide3)
    add_slide_header(slide3, "Lumenova's Milestone-Gated Escrow", "Our Solution")

    desc_box = slide3.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.3), Inches(0.6))
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    desc_p = desc_tf.paragraphs[0]
    desc_p.text = "Lumenova shifts the power balance back to the community by introducing a 4-milestone smart escrow contract. Funds are unlocked in 25% tranches only after cryptographic proof is verified and approved on-chain by the donors."
    desc_p.font.name = 'Segoe UI'
    desc_p.font.size = Pt(14)
    desc_p.font.color.rgb = TEXT_MUTED

    # 4 horizontal cards representing the flow
    step_width = Inches(2.6)
    step_height = Inches(4.0)
    gap = Inches(0.3)
    left_start = Inches(1.0)
    top_pos = Inches(2.2)

    steps = [
        ("01. Donate & Lock", "Funds are securely locked in the Soroban smart contract escrow. Donors receive Soulbound badge tokens representing governance rights.", ACCENT_CYAN),
        ("02. Reach & Submit", "Creator builds the project and completes a milestone. Creator submits progress proof (hash/IPFS link) to the smart contract.", TEXT_MAIN),
        ("03. Weighted Vote", "Donors vote to approve or reject the submitted proof. Votes are dynamically weighted based on individual donation amounts.", ACCENT_GOLD),
        ("04. Release / Refund", "Approved: 25% of funds are released to creator. Rejected: Donors can instantly claim a refund for their remaining 25% share.", ACCENT_CYAN)
    ]

    for idx, (title, body, accent_color) in enumerate(steps):
        left = left_start + idx * (step_width + gap)
        # Card background
        add_card(slide3, left, top_pos, step_width, step_height)
        
        # Accent indicator
        indicator = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top_pos, step_width, Inches(0.1))
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = accent_color
        indicator.line.fill.background()

        # Step Text
        add_text_box(slide3, left + Inches(0.2), top_pos + Inches(0.4), step_width - Inches(0.4), Inches(0.6), title, 16, True, accent_color)
        add_text_box(slide3, left + Inches(0.2), top_pos + Inches(1.2), step_width - Inches(0.4), Inches(2.5), body, 13, False, TEXT_MUTED)

        # Arrow indicator between cards
        if idx < 3:
            arrow = slide3.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, 
                left + step_width + Inches(0.05), 
                top_pos + Inches(1.8), 
                Inches(0.2), 
                Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = CARD_BORDER
            arrow.line.fill.background()

    # ----------------------------------------------------
    # SLIDE 4: How It Works (Product Walkthrough Placeholders)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_background(slide4)
    add_slide_header(slide4, "Intuitive Web3 Onboarding & Interface", "Product Walkthrough")

    # 4 columns for product views
    col_width = Inches(2.6)
    col_height = Inches(4.6)
    col_gap = Inches(0.3)
    col_left_start = Inches(1.0)
    col_top = Inches(1.8)

    walkthroughs = [
        ("1. Wallet Onboarding", "Frictionless connection using Freighter Wallet. Integrated with Stellar's Testnet Friendbot for automated, one-click test funds.", "image-2.png"),
        ("2. Interactive Grid", "Visually displays active campaigns, milestone statuses, and escrow progress. Updates dynamically via Soroban event streaming.", "image.png"),
        ("3. Trustless Escrow", "Real-time copy-to-clipboard buttons, instant status toasts, and dynamic SVG progress gauges indicating target completion rates.", "image-1.png"),
        ("4. Community Voting", "Backers can cast weighted votes. If a milestone is rejected, they can claim a refund via the refund() transaction on the contract.", "image.png")
    ]

    for idx, (title, body, img_placeholder) in enumerate(walkthroughs):
        left = col_left_start + idx * (col_width + col_gap)
        add_card(slide4, left, col_top, col_width, col_height)
        
        # Header inside card
        add_text_box(slide4, left + Inches(0.15), col_top + Inches(0.2), col_width - Inches(0.3), Inches(0.5), title, 15, True, ACCENT_GOLD)
        
        # Image placeholder box
        img_box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), col_top + Inches(0.8), col_width - Inches(0.4), Inches(1.8))
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = RGBColor(12, 17, 30)
        img_box.line.color.rgb = CARD_BORDER
        
        # Text inside image placeholder
        tx_box = slide4.shapes.add_textbox(left + Inches(0.2), col_top + Inches(1.2), col_width - Inches(0.4), Inches(1.0))
        tf_tx = tx_box.text_frame
        tf_tx.word_wrap = True
        p_tx = tf_tx.paragraphs[0]
        p_tx.text = f"[ SCREENSHOT: {img_placeholder} ]"
        p_tx.font.name = 'Segoe UI'
        p_tx.font.size = Pt(10)
        p_tx.font.bold = True
        p_tx.font.color.rgb = ACCENT_CYAN
        p_tx.alignment = PP_ALIGN.CENTER

        # Body Text
        add_text_box(slide4, left + Inches(0.15), col_top + Inches(2.8), col_width - Inches(0.3), Inches(1.6), body, 12, False, TEXT_MUTED)

    # ----------------------------------------------------
    # SLIDE 5: Market Opportunity
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_background(slide5)
    add_slide_header(slide5, "Capturing the Web3 & Social Good Market", "Market Opportunity")

    # Left Column: Market
    add_card(slide5, Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.5))
    add_text_box(slide5, Inches(1.4), Inches(2.1), Inches(4.4), Inches(0.5), "Market Size & Need", 20, True, ACCENT_CYAN)

    box_mkt = slide5.shapes.add_textbox(Inches(1.4), Inches(2.7), Inches(4.4), Inches(3.2))
    tf_mkt = box_mkt.text_frame
    tf_mkt.word_wrap = True
    
    p = tf_mkt.paragraphs[0]
    p.text = "• $25B+ Crowdfunding Market:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_mkt.add_paragraph()
    p.text = "The global market is expanding, yet transaction costs, delays, and lack of transparency remain major pain points."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)
    
    p = tf_mkt.add_paragraph()
    p.text = "• Trust Deficit:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_mkt.add_paragraph()
    p.text = "Over 12% of crowdfunding campaigns face delays or fail to deliver entirely, dampening donor confidence."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = tf_mkt.add_paragraph()
    p.text = "• Decentralized Giving (Web3):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_mkt.add_paragraph()
    p.text = "Decentralized donation volume has surged by 400% over the last three years, demanding trustless solutions."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED

    # Right Column: Why Stellar?
    add_card(slide5, Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.5))
    add_text_box(slide5, Inches(7.5), Inches(2.1), Inches(4.4), Inches(0.5), "Why Build on Stellar?", 20, True, ACCENT_GOLD)

    box_stl = slide5.shapes.add_textbox(Inches(7.5), Inches(2.7), Inches(4.4), Inches(3.2))
    tf_stl = box_stl.text_frame
    tf_stl.word_wrap = True

    p = tf_stl.paragraphs[0]
    p.text = "• Low Cost & Fast Finality:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_stl.add_paragraph()
    p.text = "Stellar ensures micro-donations are economically viable, with sub-cent network transaction fees."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)
    
    p = tf_stl.add_paragraph()
    p.text = "• Interoperability & Anchors:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_stl.add_paragraph()
    p.text = "Simplifies local currency on/off ramping, allowing non-crypto users to participate globally."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = tf_stl.add_paragraph()
    p.text = "• Soroban Rust Smart Contracts:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_stl.add_paragraph()
    p.text = "Safe, execution-optimized WASM runtime for trustless, multi-contract escrow and voting logic."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 6: Technical Architecture
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_background(slide6)
    add_slide_header(slide6, "Decoupled & Scalable Smart Contract Architecture", "Technical Architecture")

    # Columns / Layers
    layer_width = Inches(2.1)
    layer_height = Inches(4.2)
    layer_gap = Inches(0.18)
    layer_left_start = Inches(1.0)
    layer_top = Inches(2.0)

    layers = [
        ("1. Frontend Layer", "React / Vite / TypeScript\n\n• Fetches contract state via RPC.\n• Custom SVG progress widgets.\n• Real-time state-based toast notifications.", ACCENT_CYAN),
        ("2. Wallet / SDK", "Freighter Wallet & SDK\n\n• Secure transaction signature.\n• Interacts with Horizon Server & Soroban RPC endpoint.", TEXT_MAIN),
        ("3. Escrow Contract", "Crowdfunding (Soroban)\n\n• Stores milestones, goals.\n• Manages escrow balances.\n• Publishes custom events.\n• Performs voting checks.", ACCENT_GOLD),
        ("4. Badge Contract", "Rewards Badge (Soroban)\n\n• Soulbound badge tiers.\n• Multi-contract mint call.\n• Custom mint authorizations (`set_minter`).", TEXT_MAIN),
        ("5. Ledger Layer", "Stellar Testnet Network\n\n• finality under 5s.\n• Low fee processing.\n• Immutable ledger records.", ACCENT_CYAN)
    ]

    for idx, (title, body, accent_color) in enumerate(layers):
        left = layer_left_start + idx * (layer_width + layer_gap)
        add_card(slide6, left, layer_top, layer_width, layer_height)

        # Highlight header
        add_text_box(slide6, left + Inches(0.15), layer_top + Inches(0.3), layer_width - Inches(0.3), Inches(0.6), title, 14, True, accent_color)
        # Content
        add_text_box(slide6, left + Inches(0.15), layer_top + Inches(1.0), layer_width - Inches(0.3), Inches(3.0), body, 11, False, TEXT_MUTED)

        # Connection line visual
        if idx < 4:
            line = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + layer_width + Inches(0.04), layer_top + Inches(2.0), Inches(0.1), Inches(0.2))
            line.fill.solid()
            line.fill.fore_color.rgb = CARD_BORDER
            line.line.fill.background()

    # ----------------------------------------------------
    # SLIDE 7: Traction & Growth (Stat Highlights)
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_background(slide7)
    add_slide_header(slide7, "Early Testnet Traction & Growth", "Traction & Metrics")

    # Big Stats Row
    stat_width = Inches(3.4)
    stat_height = Inches(2.0)
    stat_gap = Inches(0.5)
    stat_left_start = Inches(1.0)
    stat_top = Inches(1.8)

    stats = [
        ("87+", "Active Wallets", "Onboarded testnet users participating in campaigns"),
        ("83", "On-Chain Donations", "Successful crowdfunding escrow transactions"),
        ("100%", "Verification Rate", "On-chain voting success and validation of proofs")
    ]

    for idx, (value, label, desc) in enumerate(stats):
        left = stat_left_start + idx * (stat_width + stat_gap)
        add_card(slide7, left, stat_top, stat_width, stat_height)
        
        # Big number
        num_box = add_text_box(slide7, left + Inches(0.2), stat_top + Inches(0.2), stat_width - Inches(0.4), Inches(0.7), value, 40, True, ACCENT_GOLD, PP_ALIGN.CENTER)
        # Label
        add_text_box(slide7, left + Inches(0.2), stat_top + Inches(0.9), stat_width - Inches(0.4), Inches(0.3), label, 14, True, TEXT_MAIN, PP_ALIGN.CENTER)
        # Description
        add_text_box(slide7, left + Inches(0.2), stat_top + Inches(1.2), stat_width - Inches(0.4), Inches(0.6), desc, 11, False, TEXT_MUTED, PP_ALIGN.CENTER)

    # Bottom Callout Card - User Feedback
    add_card(slide7, Inches(1.0), Inches(4.3), Inches(11.3), Inches(2.0))
    add_text_box(slide7, Inches(1.3), Inches(4.6), Inches(10.7), Inches(0.4), "Community-Driven Product Development", 16, True, ACCENT_CYAN)
    
    fb_box = slide7.shapes.add_textbox(Inches(1.3), Inches(5.0), Inches(10.7), Inches(1.0))
    fb_tf = fb_box.text_frame
    fb_tf.word_wrap = True
    p_fb = fb_tf.paragraphs[0]
    p_fb.text = "• Leveraged active feedback from 83 early testnet community contributors to iteratively refine the application.\n• Implemented 8 core community-requested features: copy-to-clipboard functionality, real-time toast alerts, Light/Dark mode options, direct onboarding faucet access, and custom visual empty states."
    p_fb.font.name = 'Segoe UI'
    p_fb.font.size = Pt(12)
    p_fb.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 8: Product Iteration (Before vs After)
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_background(slide8)
    add_slide_header(slide8, "Community Feedback & Iteration Cycles", "Product Development")

    # 3 Horizontal Cards mapping key before/after iterations
    card_w = Inches(3.6)
    card_h = Inches(4.5)
    card_g = Inches(0.25)
    card_l_start = Inches(1.0)
    card_t = Inches(1.8)

    iterations = [
        ("UX & ACCESSIBILITY", "Before:\nBackers struggled to copy wallet addresses and contract IDs manually.\n\nAfter:\nAdded copy-to-clipboard buttons and state-based toast notifications for success/failure.", ACCENT_CYAN),
        ("ONBOARDING FLOW", "Before:\nNo direct mechanism for testnet token funding for new users.\n\nAfter:\nIntegrated an onboarding helper faucet button into the welcome modal for immediate setup.", ACCENT_GOLD),
        ("THEMING & ACCENTS", "Before:\nRigid dark-only design with plain visual components.\n\nAfter:\nImplemented a system-wide Dark/Light mode toggle, customized styling, and SVG state loaders.", ACCENT_CYAN)
    ]

    for idx, (title, text, accent) in enumerate(iterations):
        left = card_l_start + idx * (card_w + card_g)
        add_card(slide8, left, card_t, card_w, card_h)
        
        # Indicator line
        line = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, card_t, card_w, Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = accent
        line.line.fill.background()

        add_text_box(slide8, left + Inches(0.25), card_t + Inches(0.3), card_w - Inches(0.5), Inches(0.5), title, 16, True, accent)
        
        # Styled body text to parse "Before:" and "After:" nicely
        body_box = slide8.shapes.add_textbox(left + Inches(0.25), card_t + Inches(0.9), card_w - Inches(0.5), Inches(3.4))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        
        paragraphs_text = text.split("\n\n")
        for p_idx, para in enumerate(paragraphs_text):
            p = body_tf.add_paragraph() if p_idx > 0 else body_tf.paragraphs[0]
            p.text = para
            p.font.name = 'Segoe UI'
            p.font.size = Pt(12)
            if "Before:" in para or "After:" in para:
                p.font.bold = True
                p.font.color.rgb = TEXT_MAIN
            else:
                p.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 9: Growth & Retention Strategy
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_background(slide9)
    add_slide_header(slide9, "Ecosystem Growth & Backer Retention", "Growth Strategy")

    # Column 1: Onboarding
    add_card(slide9, Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.5))
    add_text_box(slide9, Inches(1.3), Inches(2.1), Inches(4.6), Inches(0.5), "Onboarding & Distribution", 20, True, ACCENT_CYAN)
    
    box_gro_left = slide9.shapes.add_textbox(Inches(1.3), Inches(2.7), Inches(4.6), Inches(3.2))
    tf_gro_left = box_gro_left.text_frame
    tf_gro_left.word_wrap = True

    p = tf_gro_left.paragraphs[0]
    p.text = "• Organic Sharing Options:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_gro_left.add_paragraph()
    p.text = "Embedded dynamic 'Share' button that copies custom campaign links directly to the clipboard, simplifying social media distribution."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = tf_gro_left.add_paragraph()
    p.text = "• Frictionless Entry Point:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_gro_left.add_paragraph()
    p.text = "One-click Friendbot token minting on welcome modal allows Web3 novices to immediately fund demo campaigns."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED

    # Column 2: Engagement
    add_card(slide9, Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.5))
    add_text_box(slide9, Inches(7.4), Inches(2.1), Inches(4.6), Inches(0.5), "Engagement & Gamification", 20, True, ACCENT_GOLD)

    box_gro_right = slide9.shapes.add_textbox(Inches(7.4), Inches(2.7), Inches(4.6), Inches(3.2))
    tf_gro_right = box_gro_right.text_frame
    tf_gro_right.word_wrap = True

    p = tf_gro_right.paragraphs[0]
    p.text = "• Soulbound Badges (SBTs):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_gro_right.add_paragraph()
    p.text = "Donors are minted non-transferable reward badges (Spark, Glow, Supernova) directly matching their cumulative on-chain donation volume."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = tf_gro_right.add_paragraph()
    p.text = "• Proof-of-Contribution Reputation:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_gro_right.add_paragraph()
    p.text = "Badges serve as proof of reputation across the ecosystem. Future updates will lock voting capabilities in other platforms to high-tier badge holders."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 10: Roadmap (Timeline)
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_background(slide10)
    add_slide_header(slide10, "Project Development Roadmap", "Roadmap & Future Goals")

    # Horizontal Timeline Nodes
    node_w = Inches(2.6)
    node_h = Inches(3.6)
    node_g = Inches(0.3)
    node_l_start = Inches(1.0)
    node_t = Inches(2.2)

    milestones = [
        ("Phase 1: MVP", "COMPLETED\n\n• Core Soroban contracts.\n• Mock client testing.\n• Initial milestone lock logic.", ACCENT_CYAN),
        ("Phase 2: Growth", "IN PROGRESS\n\n• Stellar Testnet launch.\n• 80+ unique user interactions.\n• Community feedback and fixes.", ACCENT_GOLD),
        ("Phase 3: Security", "Q3 2026\n\n• Smart contract security audit.\n• Mainnet deploy preparation.\n• Public beta launch.", TEXT_MAIN),
        ("Phase 4: Scaling", "Q4 2026\n\n• Stellar Anchor integration.\n• Fiat onboarding.\n• SBT multi-dApp integration.", ACCENT_CYAN)
    ]

    for idx, (title, details, color) in enumerate(milestones):
        left = node_l_start + idx * (node_w + node_g)
        add_card(slide10, left, node_t, node_w, node_h)
        
        # Indicator circle
        circle = slide10.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.1), node_t - Inches(0.4), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        add_text_box(slide10, left + Inches(0.15), node_t + Inches(0.3), node_w - Inches(0.3), Inches(0.5), title, 15, True, color, PP_ALIGN.CENTER)
        
        # Details text
        details_box = slide10.shapes.add_textbox(left + Inches(0.15), node_t + Inches(0.9), node_w - Inches(0.3), Inches(2.5))
        details_tf = details_box.text_frame
        details_tf.word_wrap = True
        
        details_lines = details.split("\n\n")
        for l_idx, line in enumerate(details_lines):
            p = details_tf.add_paragraph() if l_idx > 0 else details_tf.paragraphs[0]
            p.text = line
            p.font.name = 'Segoe UI'
            p.font.size = Pt(11)
            if "COMPLETED" in line or "IN PROGRESS" in line or "Q3 2026" in line or "Q4 2026" in line:
                p.font.bold = True
                p.font.color.rgb = TEXT_MAIN
            else:
                p.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(4)

        if idx < 3:
            connector = slide10.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                left + node_w, 
                node_t - Inches(0.25), 
                node_g, 
                Inches(0.04)
            )
            connector.fill.solid()
            connector.fill.fore_color.rgb = CARD_BORDER
            connector.line.fill.background()

    # ----------------------------------------------------
    # SLIDE 11: Ask & Partnership
    # ----------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    add_background(slide11)
    add_slide_header(slide11, "Empowering Public Goods: Our Ask", "Ecosystem Ask")

    # Two big blocks
    add_card(slide11, Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.2))
    add_text_box(slide11, Inches(1.3), Inches(2.3), Inches(4.6), Inches(0.5), "Project Vision", 20, True, ACCENT_GOLD)

    box_vision = slide11.shapes.add_textbox(Inches(1.3), Inches(3.0), Inches(4.6), Inches(3.0))
    tf_vision = box_vision.text_frame
    tf_vision.word_wrap = True
    
    p = tf_vision.paragraphs[0]
    p.text = "• Open Source Infrastructure:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_vision.add_paragraph()
    p.text = "Lumenova is built as a public good to secure crowdfunding transparently. The code is 100% open-source."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = tf_vision.add_paragraph()
    p.text = "• Community First:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_vision.add_paragraph()
    p.text = "We focus on aligning incentives between project builders and micro-donors, minimizing Web3 onboarding friction."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED

    add_card(slide11, Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.2))
    add_text_box(slide11, Inches(7.4), Inches(2.3), Inches(4.6), Inches(0.5), "Ecosystem Support Needed", 20, True, ACCENT_CYAN)

    box_ask = slide11.shapes.add_textbox(Inches(7.4), Inches(3.0), Inches(4.6), Inches(3.0))
    tf_ask = box_ask.text_frame
    tf_ask.word_wrap = True

    p = tf_ask.paragraphs[0]
    p.text = "• Smart Contract Auditing:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    
    p = tf_ask.add_paragraph()
    p.text = "Seeking partnerships with auditing teams to review escrow and badge contract security prior to Mainnet deployment."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

    p = tf_ask.add_paragraph()
    p.text = "• Anchor Integration & Pilot Projects:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_MUTED
    
    p = tf_ask.add_paragraph()
    p.text = "Looking for partners to pilot test the milestone-escrow framework for real-world charitable giving and local grants."
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 12: Closing / Contact Slide
    # ----------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    add_background(slide12)

    # Centralized bold call to action
    cta_box = slide12.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(2.0))
    cta_tf = cta_box.text_frame
    cta_tf.word_wrap = True
    
    p_cta_title = cta_tf.paragraphs[0]
    p_cta_title.text = "Empower Trust. Build Together."
    p_cta_title.font.name = 'Trebuchet MS'
    p_cta_title.font.size = Pt(44)
    p_cta_title.font.bold = True
    p_cta_title.font.color.rgb = ACCENT_GOLD
    p_cta_title.alignment = PP_ALIGN.CENTER
    
    p_cta_desc = cta_tf.add_paragraph()
    p_cta_desc.text = "Bringing transparency and community-led alignment to global crowdfunding."
    p_cta_desc.font.name = 'Segoe UI'
    p_cta_desc.font.size = Pt(18)
    p_cta_desc.font.color.rgb = TEXT_MUTED
    p_cta_desc.alignment = PP_ALIGN.CENTER
    p_cta_desc.space_before = Pt(15)

    # 3 contact cards
    contact_w = Inches(3.4)
    contact_h = Inches(2.2)
    contact_g = Inches(0.5)
    contact_l = Inches(1.0)
    contact_t = Inches(4.0)

    contacts = [
        ("GitHub Repository", "github.com/Abhishek86038/Lumenova-A5\n\n[ Check code and commits ]", ACCENT_CYAN),
        ("Live Demonstration", "lumenova-a4.vercel.app\n\n[ Access on Testnet ]", ACCENT_GOLD),
        ("Developer Contact", "Abhishek / Team\n[ abhishek@example.com ]\n\n[ Say hello! ]", ACCENT_CYAN)
    ]

    for idx, (title, info, accent) in enumerate(contacts):
        left = contact_l + idx * (contact_w + contact_g)
        add_card(slide12, left, contact_t, contact_w, contact_h)
        
        # Header
        add_text_box(slide12, left + Inches(0.15), contact_t + Inches(0.2), contact_w - Inches(0.3), Inches(0.4), title, 15, True, accent, PP_ALIGN.CENTER)
        # Info
        add_text_box(slide12, left + Inches(0.15), contact_t + Inches(0.7), contact_w - Inches(0.3), Inches(1.2), info, 12, False, TEXT_MUTED, PP_ALIGN.CENTER)

    prs.save("Lumenova_Pitch_Deck.pptx")
    print("Pitch deck successfully created: Lumenova_Pitch_Deck.pptx")

if __name__ == '__main__':
    create_deck()
